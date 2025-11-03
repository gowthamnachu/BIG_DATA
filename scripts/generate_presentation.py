"""
Generate PowerPoint Presentation for Fraud Detection Project
Creates a comprehensive PPT with all project details, figures, and results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def add_content_slide(prs, title, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        
def add_two_column_slide(prs, title, left_content, right_content):
    """Add slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(10)

def add_image_slide(prs, title, image_path, caption=""):
    """Add slide with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Add image if exists
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.2), width=Inches(8))
        print(f"  ✓ Added image: {image_path}")
    else:
        # Add placeholder text
        text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        text_frame = text_box.text_frame
        text_frame.text = f"[Image: {os.path.basename(image_path)}]"
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.size = Pt(24)
        text_frame.paragraphs[0].font.italic = True
        print(f"  ⚠ Image not found: {image_path}")
    
    # Add caption
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        caption_frame.paragraphs[0].font.size = Pt(14)
        caption_frame.paragraphs[0].font.italic = True

def add_table_slide(prs, title, headers, rows):
    """Add slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Add table
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.5)
    
    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    
    # Set column widths
    for i in range(cols_count):
        table.columns[i].width = Inches(9.0 / cols_count)
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.rows[i + 1].cells[j]
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

# Create presentation
print("Starting PowerPoint generation...")
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
print("\n[1/20] Creating Title Slide...")
add_title_slide(prs, 
    "Fraud Detection using Big Data Analytics in Banking",
    "Machine Learning Approach with Random Forest Classifier\n\nSubmitted by:\n2310080023 : N. Gowtham\n2310080029 : M. Snighda Sharma\n2310080037 : I Sriram\n2310080038 : P. Manmith Reddy\n\nDate: November 3, 2025")

# Slide 2: Agenda
print("[2/20] Creating Agenda Slide...")
add_content_slide(prs, "Agenda", [
    "Introduction & Problem Statement",
    "Literature Review",
    "System Architecture & Methodology",
    "Dataset & Preprocessing",
    "Model Architecture & Training",
    "Experimental Results",
    "Performance Analysis",
    "Web Application Demo",
    "Conclusion & Future Work"
])

# Slide 3: Introduction
print("[3/20] Creating Introduction Slide...")
add_content_slide(prs, "Introduction", [
    "Financial fraud costs consumers $5.8 billion annually (FTC, 2021)",
    "Traditional rule-based systems struggle with evolving fraud tactics",
    "Machine learning offers powerful pattern recognition capabilities",
    "Challenge: Highly imbalanced datasets (<2% fraud transactions)",
    "Goal: Real-time fraud detection with minimal false positives"
])

# Slide 4: Problem Statement
print("[4/20] Creating Problem Statement Slide...")
add_content_slide(prs, "Problem Statement", [
    "Detect fraudulent banking transactions in real-time",
    "Handle severe class imbalance (85% normal, 15% fraud)",
    "Minimize false positives to avoid customer inconvenience",
    "Achieve zero false negatives for maximum protection",
    "Provide explainable predictions for fraud analysts",
    "Process 400-600 transactions per second for scalability"
])

# Slide 5: Literature Review
print("[5/20] Creating Literature Review Slide...")
add_two_column_slide(prs, "Literature Review",
    ["Traditional Approaches:",
     "• Rule-based systems (Phua et al., 2010)",
     "• Data mining techniques (Ngai et al., 2011)",
     "",
     "Machine Learning:",
     "• Ensemble methods (Gadi et al., 2008)",
     "• Cost-sensitive learning (Sahin & Duman, 2011)",
     "",
     "Imbalanced Data:",
     "• SMOTE technique (Chawla et al., 2002)",
     "• Undersampling (Dal Pozzolo et al., 2015)"],
    ["Deep Learning:",
     "• Autoencoders (Wang et al., 2018)",
     "• Graph-based approaches (Pourhabibi et al., 2020)",
     "",
     "Explainable AI:",
     "• SHAP values (Lundberg & Lee, 2017)",
     "• LIME (Ribeiro et al., 2016)",
     "",
     "Research Gap:",
     "• Lack of production-ready systems",
     "• Limited user-friendly interfaces",
     "• Missing real-time capabilities"])

# Slide 6: System Architecture
print("[6/20] Creating System Architecture Slide...")
add_content_slide(prs, "System Architecture", [
    "Modular Design with 5 Components:",
    "1. Data Preprocessing - Missing value handling, scaling, alignment",
    "2. Feature Engineering - 29 features (Amount + V1-V28)",
    "3. Model Training - RandomForest with 300 trees",
    "4. Prediction Pipeline - Real-time fraud detection",
    "5. Web Interface - Flask-based UI with visualizations"
])

# Slide 7: Training Workflow Image
print("[7/20] Creating Training Workflow Slide...")
add_image_slide(prs, "Training Workflow", "figures/figure_2_training_workflow.png",
    "End-to-end training process from data loading to model deployment")

# Slide 8: Dataset Description
print("[8/20] Creating Dataset Slide...")
add_content_slide(prs, "Dataset Description", [
    "Synthetic Dataset: 5,000 transactions (15% fraud ratio)",
    "Training Set: 4,000 transactions (80%)",
    "Test Set: 1,000 transactions (20%)",
    "Validation Set: 200 transactions (20% fraud)",
    "",
    "Features: 29 total",
    "• Amount: Transaction value in dollars",
    "• V1-V28: PCA-transformed anonymized features"
])

# Slide 9: Amount Distribution Image
print("[9/20] Creating Amount Distribution Slide...")
add_image_slide(prs, "Transaction Amount Distribution", "figures/figure_3_amount_distribution.png",
    "Clear separation between normal and fraudulent transaction amounts")

# Slide 10: Model Architecture
print("[10/20] Creating Model Architecture Slide...")
add_content_slide(prs, "Model Architecture", [
    "RandomForestClassifier Configuration:",
    "• n_estimators: 300 trees (increased from 100)",
    "• max_depth: 10 (prevents overfitting)",
    "• min_samples_split: 10",
    "• min_samples_leaf: 4",
    "• class_weight: 'balanced' (handles imbalance)",
    "• random_state: 42 (reproducibility)",
    "",
    "Preprocessing: StandardScaler (μ=0, σ=1)"
])

# Slide 11: Feature Importance Image
print("[11/20] Creating Feature Importance Slide...")
add_image_slide(prs, "Feature Importance Analysis", "figures/figure_7_feature_importance.png",
    "Amount contributes 30.69% - the most significant predictor")

# Slide 12: Model Comparison Table
print("[12/20] Creating Model Comparison Slide...")
add_table_slide(prs, "Model Performance Comparison",
    ["Model", "Accuracy", "Precision", "Recall", "F1-Score", "ROC-AUC"],
    [
        ["Logistic Regression", "94.2%", "88.5%", "82.0%", "85.1%", "0.9654"],
        ["Decision Tree", "96.8%", "92.3%", "90.7%", "91.5%", "0.9802"],
        ["Random Forest (100)", "100.0%", "100.0%", "100.0%", "100.0%", "1.0000"],
        ["Random Forest (300)", "100.0%", "100.0%", "100.0%", "100.0%", "1.0000"],
        ["XGBoost", "99.8%", "99.3%", "99.3%", "99.3%", "0.9998"]
    ])

# Slide 13: ROC Curve Image
print("[13/20] Creating ROC Curve Slide...")
add_image_slide(prs, "ROC Curve - Perfect Classification", "figures/figure_5_roc_curve.png",
    "ROC-AUC = 1.0000 indicates perfect classification performance")

# Slide 14: Probability Distribution Image
print("[14/20] Creating Probability Distribution Slide...")
add_image_slide(prs, "Probability Distribution Analysis", "figures/figure_4_probability_distribution.png",
    "Clear separation with optimal threshold at 0.650")

# Slide 15: Precision-Recall Curve Image
print("[15/20] Creating Precision-Recall Slide...")
add_image_slide(prs, "Precision-Recall Curve", "figures/figure_6_precision_recall_curve.png",
    "Perfect precision and recall across all threshold values")

# Slide 16: Results Summary
print("[16/20] Creating Results Summary Slide...")
add_content_slide(prs, "Key Results", [
    "Perfect Classification Performance:",
    "• Accuracy: 100.00%",
    "• Precision: 100.00% (zero false positives)",
    "• Recall: 100.00% (zero false negatives)",
    "• F1-Score: 100.00%",
    "• ROC-AUC: 1.0000",
    "",
    "High-Value Fraud Detection: 100% (>$5,000)",
    "Processing Speed: 400-600 transactions/second"
])

# Slide 17: Web Application Image
print("[17/20] Creating Web Application Slide...")
add_image_slide(prs, "Web Application Interface", "figures/figure_8_web_application_display.png",
    "Interactive visualizations with pie charts, bar charts, and detailed statistics")

# Slide 18: Confusion Matrix
print("[18/20] Creating Confusion Matrix Slide...")
add_table_slide(prs, "Confusion Matrix - Test Set (1,000 transactions)",
    ["", "Predicted Normal", "Predicted Fraud"],
    [
        ["Actual Normal", "850", "0"],
        ["Actual Fraud", "0", "150"]
    ])

# Slide 19: Conclusion
print("[19/20] Creating Conclusion Slide...")
add_content_slide(prs, "Conclusion", [
    "✓ Achieved perfect classification (100% accuracy)",
    "✓ Zero false positives - no customer inconvenience",
    "✓ Zero false negatives - maximum fraud protection",
    "✓ Real-time processing capability (400-600 txn/sec)",
    "✓ Production-ready web application with Flask",
    "✓ Explainable predictions with feature importance",
    "✓ Robust generalization across different fraud ratios"
])

# Slide 20: Future Work
print("[20/20] Creating Future Work Slide...")
add_two_column_slide(prs, "Future Work & Recommendations",
    ["Technical Enhancements:",
     "• Deep learning integration (LSTM, GNN)",
     "• Online learning for concept drift",
     "• Advanced feature engineering",
     "• Ensemble heterogeneous models",
     "• SHAP/LIME for explainability",
     "",
     "Production Features:",
     "• Microservices architecture",
     "• Database integration",
     "• RESTful API development",
     "• Real-time streaming (Kafka)"],
    ["Deployment Recommendations:",
     "• Pilot testing with real data",
     "• Dynamic threshold tuning",
     "• Human-in-the-loop for borderline cases",
     "• Continuous monitoring (KPIs)",
     "• Quarterly model retraining",
     "• A/B testing vs existing systems",
     "",
     "Extensions:",
     "• Insurance fraud detection",
     "• E-commerce fraud detection",
     "• Healthcare fraud detection"])

# Save presentation
output_file = "Fraud_Detection_Presentation_Final.pptx"
prs.save(output_file)

print("\n" + "="*60)
print("✓ POWERPOINT PRESENTATION CREATED SUCCESSFULLY!")
print("="*60)
print(f"\nFile saved as: {output_file}")
print(f"Total slides: 20")
print("\nSlide breakdown:")
print("  1. Title Slide")
print("  2. Agenda")
print("  3. Introduction")
print("  4. Problem Statement")
print("  5. Literature Review")
print("  6. System Architecture")
print("  7. Training Workflow (with diagram)")
print("  8. Dataset Description")
print("  9. Amount Distribution (with chart)")
print("  10. Model Architecture")
print("  11. Feature Importance (with chart)")
print("  12. Model Comparison (table)")
print("  13. ROC Curve (with chart)")
print("  14. Probability Distribution (with chart)")
print("  15. Precision-Recall Curve (with chart)")
print("  16. Results Summary")
print("  17. Web Application (with screenshot)")
print("  18. Confusion Matrix (table)")
print("  19. Conclusion")
print("  20. Future Work & Recommendations")
print("="*60)
